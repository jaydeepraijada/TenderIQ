#!/usr/bin/env python3
"""TenderIQ v1 — Dark Professional (PPTX)"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Palette
BG    = RGBColor(0x0D, 0x1B, 0x2A)
CARD  = RGBColor(0x1E, 0x3A, 0x5F)
CARD2 = RGBColor(0x16, 0x2D, 0x4A)
BDR   = RGBColor(0x2D, 0x4A, 0x6B)
GOLD  = RGBColor(0xF0, 0xA5, 0x00)
WHT   = RGBColor(0xF1, 0xF5, 0xF9)
MUT   = RGBColor(0x94, 0xA3, 0xB8)
GRN   = RGBColor(0x22, 0xC5, 0x5E)
RED_C = RGBColor(0xEF, 0x44, 0x44)
AMB   = RGBColor(0xF5, 0x9E, 0x0B)
FONT  = "Arial"


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h, sz=15, bold=False, clr=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.name = FONT
    if clr:
        r.font.color.rgb = clr
    return tb


def heading(slide, text, y=Inches(0.3)):
    rect(slide, Inches(0.5), y + Inches(0.05), Inches(0.07), Inches(0.48), GOLD)
    txt(slide, text, Inches(0.7), y, Inches(12.3), Inches(0.6), sz=26, bold=True, clr=GOLD)


def bg(slide):
    rect(slide, 0, 0, Inches(13.33), Inches(7.5), BG)


def card_box(slide, x, y, w, h, title=None, lines=None, title_clr=None, accent=None):
    border_clr = accent or GOLD
    rect(slide, x, y, w, h, CARD, border_clr, Pt(1))
    cy = y + Inches(0.18)
    if title:
        txt(slide, title, x + Inches(0.18), cy, w - Inches(0.36), Inches(0.45),
            sz=15, bold=True, clr=title_clr or GOLD)
        cy += Inches(0.45)
    if lines:
        for line in lines:
            txt(slide, line, x + Inches(0.18), cy, w - Inches(0.36), Inches(0.35),
                sz=12, clr=WHT)
            cy += Inches(0.32)


def chip(slide, x, y, w, h, fill, text):
    rect(slide, x, y, w, h, fill)
    txt(slide, text, x, y, w, h, sz=12, bold=True,
        clr=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)


# ── Slide 1 ─────────────────────────────────────────────────────────────────
def s1(prs):
    slide = blank(prs)
    bg(slide)
    rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.6), RGBColor(0x1E, 0x3A, 0x5F))
    txt(slide, "⚖", Inches(5.2), Inches(0.8), Inches(3), Inches(1.6),
        sz=80, align=PP_ALIGN.CENTER, clr=GOLD)
    txt(slide, "TenderIQ", Inches(1), Inches(2.4), Inches(11.33), Inches(1.4),
        sz=64, bold=True, clr=GOLD, align=PP_ALIGN.CENTER)
    txt(slide, "Explainable AI for Government Tender Evaluation",
        Inches(1), Inches(3.85), Inches(11.33), Inches(0.65),
        sz=22, clr=WHT, align=PP_ALIGN.CENTER)
    txt(slide, "CRPF Hackathon  ·  Theme 3",
        Inches(1), Inches(4.6), Inches(11.33), Inches(0.5),
        sz=17, clr=MUT, align=PP_ALIGN.CENTER)
    txt(slide, "From days to minutes. Every decision traceable.",
        Inches(1), Inches(5.25), Inches(11.33), Inches(0.5),
        sz=16, clr=GOLD, align=PP_ALIGN.CENTER)
    rect(slide, Inches(3.5), Inches(6.85), Inches(6.33), Inches(0.05), GOLD)


# ── Slide 2 ─────────────────────────────────────────────────────────────────
def s2(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "The Problem with Manual Tender Evaluation")
    callouts = [
        ("3–5 Days", "per tender evaluation\nby committee"),
        ("Inconsistent", "two evaluators,\ntwo different conclusions"),
        ("No Audit Trail", "decisions made\nuntraceably"),
    ]
    cw, ch = Inches(3.8), Inches(2.2)
    for i, (big, sub) in enumerate(callouts):
        sx = Inches(0.55) + i * Inches(4.18)
        rect(slide, sx, Inches(1.1), cw, ch, CARD, GOLD, Pt(1))
        txt(slide, big, sx + Inches(0.15), Inches(1.3), cw - Inches(0.3), Inches(0.85),
            sz=30, bold=True, clr=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, sub, sx + Inches(0.15), Inches(2.2), cw - Inches(0.3), Inches(0.85),
            sz=14, clr=WHT, align=PP_ALIGN.CENTER)
    bullets = [
        "• Mixed document formats: typed PDFs, scans, phone photographs",
        "• Government procurement worth ₹50 lakh crore+ annually in India",
        "• Project delays traced directly to procurement bottlenecks",
    ]
    for i, b in enumerate(bullets):
        txt(slide, b, Inches(0.7), Inches(3.55) + Inches(0.43 * i),
            Inches(12.3), Inches(0.38), sz=14, clr=MUT)


# ── Slide 3 ─────────────────────────────────────────────────────────────────
def s3(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "TenderIQ — Four Stages, End to End")
    stages = [
        ("Stage 1 — Extract",
         "DeepSeek LLM reads the tender PDF\nand returns each criterion as structured\nJSON: category, mandatory flag, threshold\nrule, source clause, query hints."),
        ("Stage 2 — OCR & Index",
         "Three-tier pipeline handles any document\nformat. All text chunked and indexed for\nsemantic retrieval."),
        ("Stage 3 — Evaluate",
         "Vector search finds relevant evidence.\nLLM produces a verdict with confidence.\nSafety rule prevents silent disqualification."),
        ("Stage 4 — Review & Audit",
         "Borderline cases go to human review\nqueue. Every action logged. Full audit\ntrail exportable as CSV."),
    ]
    cw, ch = Inches(2.9), Inches(3.85)
    for i, (title, body) in enumerate(stages):
        sx = Inches(0.5) + i * Inches(3.15)
        rect(slide, sx, Inches(1.1), cw, ch, CARD, GOLD, Pt(1))
        rect(slide, sx, Inches(1.1), cw, Inches(0.07), GOLD)
        txt(slide, title, sx + Inches(0.15), Inches(1.2), cw - Inches(0.3), Inches(0.5),
            sz=15, bold=True, clr=GOLD)
        txt(slide, body, sx + Inches(0.15), Inches(1.75), cw - Inches(0.3), Inches(2.9),
            sz=13, clr=WHT)
    rect(slide, Inches(0.7), Inches(5.3), Inches(11.93), Inches(0.8),
         RGBColor(0x1A, 0x30, 0x50), GOLD, Pt(2))
    txt(slide, '"Minutes, not days. Every verdict traceable to a document and page."',
        Inches(0.9), Inches(5.42), Inches(11.5), Inches(0.6),
        sz=15, bold=True, clr=GOLD, align=PP_ALIGN.CENTER)


# ── Slide 4 ─────────────────────────────────────────────────────────────────
def s4(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "System Architecture")

    bw, bh = Inches(2.5), Inches(0.52)

    def abox(x, y, label, fill=CARD):
        rect(slide, x, y, bw, bh, fill, GOLD, Pt(1))
        txt(slide, label, x + Inches(0.1), y + Inches(0.06),
            bw - Inches(0.2), bh - Inches(0.12),
            sz=12, clr=WHT, align=PP_ALIGN.CENTER)

    def arr(x, y):
        rect(slide, x + Inches(1.15), y, Inches(0.2), Inches(0.3), GOLD)

    # Left: Tender PDF path
    lx = Inches(0.5)
    abox(lx, Inches(1.0), "Tender PDF")
    arr(lx, Inches(1.52))
    abox(lx, Inches(1.82), "DeepSeek LLM\n(Extract Criteria)")
    arr(lx, Inches(2.34))
    abox(lx, Inches(2.64), "Criteria JSON\n(C1–C5 structured)")

    # Right: Bidder docs path
    rx = Inches(4.0)
    abox(rx, Inches(1.0), "Bidder Docs\n(PDFs · scans · photos)")
    arr(rx, Inches(1.52))
    abox(rx, Inches(1.82), "3-Tier OCR Pipeline\n① PyMuPDF ② Tesseract ③ Vision LLM")
    arr(rx, Inches(2.34))
    abox(rx, Inches(2.64), "Vector Index\n(all-MiniLM-L6-v2 embeddings)")

    # Converging arrow
    cx = Inches(3.25)
    arr(cx, Inches(3.16))
    abox(Inches(2.0), Inches(3.46), "DeepSeek LLM — Evaluate each criterion")

    # Outputs
    rect(slide, Inches(0.7), Inches(4.2), Inches(2.0), Inches(0.52),
         RGBColor(0x14, 0x33, 0x14), GRN, Pt(1))
    txt(slide, "eligible /\nnot_eligible", Inches(0.7), Inches(4.2), Inches(2.0), Inches(0.52),
        sz=11, clr=GRN, align=PP_ALIGN.CENTER)

    rect(slide, Inches(3.3), Inches(4.2), Inches(2.0), Inches(0.52),
         RGBColor(0x3D, 0x28, 0x00), AMB, Pt(1))
    txt(slide, "needs_review\nHuman Review Queue", Inches(3.3), Inches(4.2), Inches(2.0), Inches(0.52),
        sz=11, clr=AMB, align=PP_ALIGN.CENTER)

    abox(Inches(1.5), Inches(5.1), "SQLite Audit Log", fill=RGBColor(0x0D, 0x20, 0x38))

    # Key facts sidebar
    rect(slide, Inches(7.0), Inches(1.0), Inches(5.8), Inches(5.2), CARD, GOLD, Pt(1))
    txt(slide, "Key Technical Facts", Inches(7.15), Inches(1.1), Inches(5.5), Inches(0.4),
        sz=15, bold=True, clr=GOLD)
    facts = [
        "Single-process Streamlit app — no separate backend",
        "Deployable to Streamlit Cloud or HuggingFace Spaces",
        "All storage is local: SQLite + in-memory vector index",
        "Only external dependency: DeepSeek API",
    ]
    for i, f in enumerate(facts):
        txt(slide, f"• {f}", Inches(7.15), Inches(1.65) + Inches(0.65 * i),
            Inches(5.5), Inches(0.55), sz=13, clr=WHT)


# ── Slide 5 ─────────────────────────────────────────────────────────────────
def s5(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Three-Tier OCR — Handling Any Document Format")
    tiers = [
        ("Tier 1 — PyMuPDF",
         "Trigger: Typed / digital PDF\n"
         "Cost: Free, instant\n"
         "Confidence: 1.0 (lossless text extraction)\n"
         "Source label: Typed PDF",
         RGBColor(0x1D, 0x4E, 0xD8)),
        ("Tier 2 — Tesseract",
         "Trigger: Scanned PDF or image file\n"
         "Cost: Free, local, fast\n"
         "Confidence: Mean of per-word OCR scores\n"
         "Source label: Tesseract",
         RGBColor(0x7C, 0x3A, 0xED)),
        ("Tier 3 — DeepSeek Vision LLM",
         "Trigger: Tesseract confidence < 65%\n"
         "Cost: One API call\n"
         "Confidence: 0.95\n"
         "Source label: Vision LLM\n"
         "Logged: vision_ocr_invoked",
         RGBColor(0xC2, 0x51, 0x0E)),
    ]
    cw, ch = Inches(3.6), Inches(3.2)
    for i, (title, body, accent) in enumerate(tiers):
        sx = Inches(0.5) + i * Inches(4.3)
        rect(slide, sx, Inches(1.1), cw, ch, CARD, accent, Pt(2))
        rect(slide, sx, Inches(1.1), cw, Inches(0.1), accent)
        txt(slide, title, sx + Inches(0.15), Inches(1.25), cw - Inches(0.3), Inches(0.5),
            sz=15, bold=True, clr=accent)
        txt(slide, body, sx + Inches(0.15), Inches(1.85), cw - Inches(0.3), Inches(2.2),
            sz=13, clr=WHT)

    rect(slide, Inches(0.5), Inches(4.55), Inches(12.33), Inches(1.72),
         RGBColor(0x1A, 0x28, 0x40), AMB, Pt(2))
    txt(slide, "Demo Scenario",
        Inches(0.7), Inches(4.65), Inches(12), Inches(0.38),
        sz=14, bold=True, clr=AMB)
    demo = ("Bidder C submits a blurry, rotated CA certificate scan. "
            "Tesseract reads it at ~55% confidence. Vision LLM transcribes the turnover figure correctly. "
            "Combined confidence = 0.58 → routed to human review. "
            "This is intentional — borderline evidence requires a human.")
    txt(slide, demo, Inches(0.7), Inches(5.1), Inches(12), Inches(1.0), sz=13, clr=WHT)


# ── Slide 6 ─────────────────────────────────────────────────────────────────
def s6(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Every Decision is Explainable and Auditable")

    # Left
    lx = Inches(0.5)
    rect(slide, lx, Inches(1.1), Inches(5.85), Inches(4.0), CARD, GOLD, Pt(1))
    txt(slide, "Criterion-Level Verdicts", lx + Inches(0.15), Inches(1.2),
        Inches(5.5), Inches(0.42), sz=15, bold=True, clr=GOLD)
    left_lines = [
        "Each (bidder × criterion) pair shows:",
        "• Which criterion was checked",
        "• Which document and page provided evidence",
        "• What value was extracted (e.g. 'INR 6.2 Cr')",
        "• Which OCR tier read the document",
        "• Combined confidence score (0–100%)",
        "• Plain-English reason",
    ]
    for i, line in enumerate(left_lines):
        clr = MUT if i == 0 else WHT
        txt(slide, line, lx + Inches(0.15), Inches(1.65) + Inches(0.42 * i),
            Inches(5.55), Inches(0.38), sz=13, clr=clr)

    # Right
    rx = Inches(6.98)
    rect(slide, rx, Inches(1.1), Inches(5.85), Inches(4.0), CARD, GOLD, Pt(1))
    txt(slide, "Audit Trail", rx + Inches(0.15), Inches(1.2),
        Inches(5.5), Inches(0.42), sz=15, bold=True, clr=GOLD)
    right_lines = [
        "Every action logged with:",
        "• UTC timestamp",
        "• Action type: criteria_extracted / bidder_processed",
        "  criterion_evaluated / human_review_action",
        "  vision_ocr_invoked / precomputed_fallback_used",
        "• Model version & Actor (system / officer)",
        "• Full payload JSON — Exportable as CSV",
    ]
    for i, line in enumerate(right_lines):
        clr = MUT if i == 0 else WHT
        txt(slide, line, rx + Inches(0.15), Inches(1.65) + Inches(0.42 * i),
            Inches(5.55), Inches(0.38), sz=13, clr=clr)

    # Safety rule
    rect(slide, Inches(0.5), Inches(5.38), Inches(12.33), Inches(1.22),
         RGBColor(0x3D, 0x1A, 0x00), AMB, Pt(2))
    txt(slide, "The Safety Rule:",
        Inches(0.7), Inches(5.45), Inches(4.5), Inches(0.4),
        sz=15, bold=True, clr=AMB)
    txt(slide, ("If combined confidence is 0.55–0.80 AND verdict is not_eligible, "
                "the verdict is automatically downgraded to needs_review. "
                "A bidder is NEVER silently disqualified at medium confidence."),
        Inches(0.7), Inches(5.88), Inches(12), Inches(0.65),
        sz=13, clr=WHT)


# ── Slide 7 ─────────────────────────────────────────────────────────────────
def s7(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Demo: Three Bidders, Three Outcomes")
    bidders = [
        ("Bidder A — Apex Constructions Pvt. Ltd.", GRN, "ELIGIBLE",
         ["C1 Turnover: INR 6.37 Cr avg (threshold 5 Cr) — PASS",
          "C2 Projects: 5 completed incl. CRPF barracks — PASS",
          "C3 GST: GSTIN 27AABCA1234F1Z5, Active — PASS",
          "C4 ISO 9001:2015: Valid June 2027 — PASS",
          "All typed PDFs, confidence ≥ 93% on all criteria"]),
        ("Bidder B — BuildRight Enterprises", RED_C, "NOT ELIGIBLE",
         ["C1 Turnover: INR 1.5 Cr avg (threshold 5 Cr) — FAIL",
          "\"INR 1.5 Cr is below required minimum of INR 5 Cr\"",
          "C2–C4: All pass",
          "Auto-disqualified at high confidence (95%)"]),
        ("Bidder C — Shree Constructions & Services", AMB, "NEEDS REVIEW",
         ["C1 Turnover: Submitted as blurry scan",
          "Tesseract ~55% → Vision LLM: INR 5.4 Cr",
          "Combined confidence 0.58 → safety rule triggered",
          "C2: Exactly 3 projects (borderline)",
          "C3–C4: Pass"]),
    ]
    cw, ch = Inches(3.85), Inches(3.6)
    for i, (name, color, verdict, bullets) in enumerate(bidders):
        sx = Inches(0.5) + i * Inches(4.25)
        rect(slide, sx, Inches(1.1), cw, ch, CARD, color, Pt(2))
        txt(slide, name, sx + Inches(0.15), Inches(1.2), cw - Inches(0.3), Inches(0.5),
            sz=12, bold=True, clr=WHT)
        chip(slide, sx + Inches(0.15), Inches(1.75), Inches(3.3), Inches(0.35), color, verdict)
        for j, b in enumerate(bullets):
            txt(slide, b, sx + Inches(0.15), Inches(2.2) + Inches(0.37 * j),
                cw - Inches(0.3), Inches(0.33), sz=11, clr=WHT)

    # Metric strip
    metrics = [
        ("Criteria extracted", "5"),
        ("Bidder docs processed", "15"),
        ("LLM evaluation calls", "15"),
        ("Vision OCR invocations", "1"),
        ("Human review items", "1"),
        ("Total audit entries", "20+"),
    ]
    rect(slide, Inches(0.5), Inches(5.0), Inches(12.33), Inches(0.98),
         RGBColor(0x10, 0x26, 0x40), GOLD, Pt(1))
    mw = Inches(2.0)
    for i, (label, val) in enumerate(metrics):
        mx = Inches(0.65) + i * Inches(2.0)
        txt(slide, val, mx, Inches(5.07), mw, Inches(0.38),
            sz=20, bold=True, clr=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, label, mx, Inches(5.5), mw, Inches(0.35),
            sz=11, clr=MUT, align=PP_ALIGN.CENTER)


# ── Slide 8 ─────────────────────────────────────────────────────────────────
def s8(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Stack, Impact & What's Next")
    stack = [
        ("UI & orchestration", "Streamlit 1.39"),
        ("LLM", "DeepSeek API (OpenAI-compatible)"),
        ("OCR Tier 1", "PyMuPDF 1.24"),
        ("OCR Tier 2", "Tesseract"),
        ("OCR Tier 3", "DeepSeek Vision LLM"),
        ("Semantic retrieval", "sentence-transformers all-MiniLM-L6-v2"),
        ("Data validation", "Pydantic v2"),
        ("Audit log", "SQLite"),
        ("Deployment", "Streamlit Cloud / HuggingFace Spaces"),
    ]
    rect(slide, Inches(0.5), Inches(1.05), Inches(6.2), Inches(0.42),
         RGBColor(0x0D, 0x28, 0x48), GOLD, Pt(1))
    txt(slide, "Component", Inches(0.6), Inches(1.1), Inches(2.8), Inches(0.32),
        sz=13, bold=True, clr=GOLD)
    txt(slide, "Technology", Inches(3.5), Inches(1.1), Inches(3.0), Inches(0.32),
        sz=13, bold=True, clr=GOLD)
    for i, (comp, tech) in enumerate(stack):
        row = CARD if i % 2 == 0 else CARD2
        ry = Inches(1.47) + Inches(0.39 * i)
        rect(slide, Inches(0.5), ry, Inches(6.2), Inches(0.39), row)
        txt(slide, comp, Inches(0.6), ry + Inches(0.06), Inches(2.8), Inches(0.3),
            sz=12, clr=MUT)
        txt(slide, tech, Inches(3.5), ry + Inches(0.06), Inches(3.0), Inches(0.3),
            sz=12, clr=WHT)

    txt(slide, "What's Next", Inches(7.3), Inches(1.05), Inches(5.6), Inches(0.42),
        sz=16, bold=True, clr=GOLD)
    future = [
        "Multi-tender workspace — same bidder pool, multiple tenders",
        "GeM portal API integration — live tender ingestion",
        "Automated bidder ranking with weighted scoring",
        "LayoutLM for complex financial tables in scanned statements",
        "Multi-evaluator workflow with role-based approval",
        "Review queue email/SMS notifications",
        "Audit PDF export for procurement oversight submissions",
    ]
    for i, f in enumerate(future):
        txt(slide, f"• {f}", Inches(7.3), Inches(1.55) + Inches(0.47 * i),
            Inches(5.8), Inches(0.42), sz=13, clr=WHT)

    rect(slide, Inches(0.5), Inches(6.15), Inches(12.33), Inches(1.0),
         RGBColor(0x1A, 0x32, 0x50), GOLD, Pt(2))
    txt(slide, ("3–5 days → minutes.  Every verdict traceable to a document, page, and model version."
                "  Built in one hackathon session. Deployable today."),
        Inches(0.7), Inches(6.25), Inches(12), Inches(0.8),
        sz=14, bold=True, clr=GOLD, align=PP_ALIGN.CENTER)


def main():
    os.makedirs("deck", exist_ok=True)
    prs = new_prs()
    s1(prs); s2(prs); s3(prs); s4(prs)
    s5(prs); s6(prs); s7(prs); s8(prs)
    out = "deck/TenderIQ_v1_dark_professional.pptx"
    prs.save(out)
    print(f"OK: Saved {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
