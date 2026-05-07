#!/usr/bin/env python3
"""TenderIQ v2 — Clean Minimal (PPTX)"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Palette
BG    = RGBColor(0xFF, 0xFF, 0xFF)
TXT1  = RGBColor(0x11, 0x18, 0x27)
TXT2  = RGBColor(0x37, 0x41, 0x51)
TXT3  = RGBColor(0x6B, 0x72, 0x80)
ACC   = RGBColor(0x25, 0x63, 0xEB)
ACBG  = RGBColor(0xEF, 0xF6, 0xFF)
DIV   = RGBColor(0xE5, 0xE7, 0xEB)
GRN   = RGBColor(0x05, 0x96, 0x69)
GRNBG = RGBColor(0xD1, 0xFA, 0xE5)
RED_C = RGBColor(0xDC, 0x26, 0x26)
REDBG = RGBColor(0xFE, 0xE2, 0xE2)
AMB   = RGBColor(0xD9, 0x77, 0x06)
AMBBG = RGBColor(0xFE, 0xF3, 0xC7)
FONT  = "Calibri"


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h, sz=15, bold=False, clr=None, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = FONT
    if clr:
        r.font.color.rgb = clr
    return tb


def bg(slide):
    rect(slide, 0, 0, Inches(13.33), Inches(7.5), BG)


def heading(slide, text, y=Inches(0.35)):
    rect(slide, Inches(0.75), y + Inches(0.04), Inches(0.05), Inches(0.5), ACC)
    txt(slide, text, Inches(0.95), y, Inches(12.0), Inches(0.6),
        sz=28, bold=True, clr=TXT1)


def card(slide, x, y, w, h):
    rect(slide, x, y, w, h, RGBColor(0xFA, 0xFB, 0xFF), DIV, Pt(1))


def chip(slide, x, y, w, h, bg_c, text_c, text):
    rect(slide, x, y, w, h, bg_c)
    txt(slide, text, x, y, w, h, sz=12, bold=True, clr=text_c, align=PP_ALIGN.CENTER)


# ── Slide 1 ─────────────────────────────────────────────────────────────────
def s1(prs):
    slide = blank(prs)
    bg(slide)
    rect(slide, Inches(0.75), Inches(2.2), Inches(0.05), Inches(3.0), ACC)
    txt(slide, "⚖", Inches(1.0), Inches(2.2), Inches(1.5), Inches(1.2),
        sz=56, clr=ACC)
    txt(slide, "TenderIQ", Inches(2.7), Inches(2.2), Inches(10), Inches(1.2),
        sz=64, bold=True, clr=TXT1)
    txt(slide, "Explainable AI for Government Tender Evaluation",
        Inches(2.7), Inches(3.5), Inches(10), Inches(0.7),
        sz=22, clr=TXT3)
    rect(slide, Inches(2.7), Inches(4.32), Inches(7), Inches(0.03), DIV)
    txt(slide, "CRPF Hackathon  ·  Theme 3",
        Inches(2.7), Inches(4.5), Inches(6), Inches(0.45),
        sz=16, clr=TXT3)
    txt(slide, "From days to minutes. Every decision traceable.",
        Inches(2.7), Inches(5.05), Inches(8), Inches(0.45),
        sz=15, italic=True, clr=ACC)


# ── Slide 2 ─────────────────────────────────────────────────────────────────
def s2(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "The Problem with Manual Tender Evaluation")
    callouts = [
        ("3–5", "Days per tender evaluation"),
        ("≠", "Inconsistent verdicts\nfrom different evaluators"),
        ("0", "Zero audit trail\nfor decisions made"),
    ]
    cw = Inches(3.5)
    for i, (big, sub) in enumerate(callouts):
        sx = Inches(0.75) + i * Inches(4.15)
        card(slide, sx, Inches(1.15), cw, Inches(2.4))
        txt(slide, big, sx + Inches(0.2), Inches(1.3), cw - Inches(0.4), Inches(1.1),
            sz=72, bold=True, clr=ACC, align=PP_ALIGN.CENTER)
        txt(slide, sub, sx + Inches(0.2), Inches(2.45), cw - Inches(0.4), Inches(0.85),
            sz=14, clr=TXT2, align=PP_ALIGN.CENTER)
    rect(slide, Inches(0.75), Inches(3.8), Inches(11.83), Inches(0.03), DIV)
    bullets = [
        "• Mixed document formats: typed PDFs, scans, phone photographs",
        "• Government procurement worth ₹50 lakh crore+ annually in India",
        "• Project delays traced directly to procurement bottlenecks",
    ]
    for i, b in enumerate(bullets):
        txt(slide, b, Inches(0.95), Inches(4.0) + Inches(0.45 * i),
            Inches(12), Inches(0.38), sz=14, clr=TXT3)


# ── Slide 3 ─────────────────────────────────────────────────────────────────
def s3(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "TenderIQ — Four Stages, End to End")
    stages = [
        ("1", "Extract",
         "DeepSeek LLM reads the tender PDF\nand returns each criterion as structured JSON:\ncategory, mandatory flag, threshold rule,\nsource clause, query hints."),
        ("2", "OCR & Index",
         "Three-tier pipeline handles any document\nformat. All text chunked and indexed for\nsemantic retrieval."),
        ("3", "Evaluate",
         "Vector search finds relevant evidence.\nLLM produces a verdict with confidence.\nSafety rule prevents silent disqualification."),
        ("4", "Review & Audit",
         "Borderline cases go to human review\nqueue. Every action logged. Full audit\ntrail exportable as CSV."),
    ]
    cw, ch = Inches(2.9), Inches(3.7)
    for i, (num, title, body) in enumerate(stages):
        sx = Inches(0.75) + i * Inches(3.15)
        card(slide, sx, Inches(1.1), cw, ch)
        rect(slide, sx, Inches(1.1), Inches(0.05), Inches(0.6), ACC)
        txt(slide, num, sx + Inches(0.12), Inches(1.12), Inches(0.5), Inches(0.55),
            sz=28, bold=True, clr=ACC)
        txt(slide, title, sx + Inches(0.7), Inches(1.15), cw - Inches(0.85), Inches(0.5),
            sz=17, bold=True, clr=TXT1)
        txt(slide, body, sx + Inches(0.12), Inches(1.75), cw - Inches(0.25), Inches(2.8),
            sz=13, clr=TXT2)

    rect(slide, Inches(0.75), Inches(5.1), Inches(11.83), Inches(0.85),
         ACBG, ACC, Pt(1))
    txt(slide, '"Minutes, not days. Every verdict traceable to a document and page."',
        Inches(0.95), Inches(5.22), Inches(11.5), Inches(0.65),
        sz=15, bold=True, clr=ACC, align=PP_ALIGN.CENTER)


# ── Slide 4 ─────────────────────────────────────────────────────────────────
def s4(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "System Architecture")

    bw, bh = Inches(2.5), Inches(0.52)
    grey = RGBColor(0xF3, 0xF4, 0xF6)

    def abox(x, y, label):
        rect(slide, x, y, bw, bh, grey, DIV, Pt(1))
        txt(slide, label, x + Inches(0.1), y + Inches(0.06),
            bw - Inches(0.2), bh - Inches(0.1),
            sz=12, clr=TXT1, align=PP_ALIGN.CENTER)

    def arr(x, y):
        rect(slide, x + Inches(1.15), y, Inches(0.2), Inches(0.28), ACC)

    lx = Inches(0.75)
    abox(lx, Inches(1.0), "Tender PDF")
    arr(lx, Inches(1.52))
    abox(lx, Inches(1.8), "DeepSeek LLM\n(Extract Criteria)")
    arr(lx, Inches(2.32))
    abox(lx, Inches(2.6), "Criteria JSON\n(C1–C5 structured)")

    rx = Inches(4.2)
    abox(rx, Inches(1.0), "Bidder Docs\n(PDFs · scans · photos)")
    arr(rx, Inches(1.52))
    abox(rx, Inches(1.8), "3-Tier OCR\n① PyMuPDF ② Tesseract ③ Vision LLM")
    arr(rx, Inches(2.32))
    abox(rx, Inches(2.6), "Vector Index\n(all-MiniLM-L6-v2 embeddings)")

    cx = Inches(3.25)
    arr(cx, Inches(3.12))
    abox(Inches(2.25), Inches(3.4), "DeepSeek LLM — Evaluate each criterion")

    rect(slide, Inches(1.05), Inches(4.15), Inches(2.0), Inches(0.52),
         GRNBG, GRN, Pt(1))
    txt(slide, "eligible /\nnot_eligible", Inches(1.05), Inches(4.15), Inches(2.0), Inches(0.52),
        sz=11, clr=GRN, align=PP_ALIGN.CENTER)

    rect(slide, Inches(3.6), Inches(4.15), Inches(2.0), Inches(0.52),
         AMBBG, AMB, Pt(1))
    txt(slide, "needs_review\nHuman Review Queue", Inches(3.6), Inches(4.15), Inches(2.0), Inches(0.52),
        sz=11, clr=AMB, align=PP_ALIGN.CENTER)

    abox(Inches(2.25), Inches(5.0), "SQLite Audit Log")

    # Key facts
    rect(slide, Inches(7.3), Inches(1.0), Inches(5.6), Inches(5.0), ACBG, ACC, Pt(1))
    txt(slide, "Key Technical Facts", Inches(7.45), Inches(1.1), Inches(5.3), Inches(0.42),
        sz=15, bold=True, clr=ACC)
    facts = [
        "Single-process Streamlit app — no separate backend",
        "Deployable to Streamlit Cloud or HuggingFace Spaces",
        "All storage is local: SQLite + in-memory vector index",
        "Only external dependency: DeepSeek API",
    ]
    for i, f in enumerate(facts):
        txt(slide, f"• {f}", Inches(7.45), Inches(1.65) + Inches(0.65 * i),
            Inches(5.3), Inches(0.55), sz=13, clr=TXT2)


# ── Slide 5 ─────────────────────────────────────────────────────────────────
def s5(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Three-Tier OCR — Handling Any Document Format")
    tiers = [
        ("Tier 1", "PyMuPDF",
         "Trigger: Typed / digital PDF\nCost: Free, instant\nConfidence: 1.0 (lossless text)\nSource label: Typed PDF",
         ACC, ACBG),
        ("Tier 2", "Tesseract",
         "Trigger: Scanned PDF or image\nCost: Free, local, fast\nConfidence: Mean per-word OCR scores\nSource label: Tesseract",
         RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xED, 0xE9, 0xFE)),
        ("Tier 3", "DeepSeek Vision LLM",
         "Trigger: Tesseract confidence < 65%\nCost: One API call\nConfidence: 0.95\nSource label: Vision LLM\nLogged: vision_ocr_invoked",
         RGBColor(0xEA, 0x58, 0x0C), RGBColor(0xFF, 0xED, 0xD5)),
    ]
    cw, ch = Inches(3.6), Inches(3.2)
    for i, (tier, name, body, accent, abg) in enumerate(tiers):
        sx = Inches(0.75) + i * Inches(4.15)
        card(slide, sx, Inches(1.1), cw, ch)
        rect(slide, sx, Inches(1.1), cw, Inches(0.08), accent)
        txt(slide, tier, sx + Inches(0.15), Inches(1.23), Inches(1.2), Inches(0.38),
            sz=12, bold=True, clr=accent)
        txt(slide, name, sx + Inches(0.15), Inches(1.65), cw - Inches(0.3), Inches(0.42),
            sz=16, bold=True, clr=TXT1)
        txt(slide, body, sx + Inches(0.15), Inches(2.12), cw - Inches(0.3), Inches(2.0),
            sz=13, clr=TXT2)

    rect(slide, Inches(0.75), Inches(4.55), Inches(11.83), Inches(1.75),
         ACBG, ACC, Pt(1))
    txt(slide, "Demo Scenario",
        Inches(0.95), Inches(4.65), Inches(8), Inches(0.38),
        sz=14, bold=True, clr=ACC)
    demo = ("Bidder C submits a blurry, rotated CA certificate scan. "
            "Tesseract reads it at ~55% confidence. Vision LLM transcribes the turnover figure correctly. "
            "Combined confidence = 0.58 → routed to human review. "
            "This is intentional — borderline evidence requires a human.")
    txt(slide, demo, Inches(0.95), Inches(5.1), Inches(11.6), Inches(1.1), sz=13, clr=TXT2)


# ── Slide 6 ─────────────────────────────────────────────────────────────────
def s6(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Every Decision is Explainable and Auditable")

    lx = Inches(0.75)
    card(slide, lx, Inches(1.1), Inches(5.7), Inches(4.0))
    rect(slide, lx, Inches(1.1), Inches(0.05), Inches(4.0), ACC)
    txt(slide, "Criterion-Level Verdicts", lx + Inches(0.2), Inches(1.2),
        Inches(5.4), Inches(0.42), sz=16, bold=True, clr=TXT1)
    left_lines = [
        "Each (bidder × criterion) pair shows:",
        "• Which criterion was checked",
        "• Which document and page provided evidence",
        "• What value was extracted (e.g. 'INR 6.2 Cr')",
        "• Which OCR tier read the document",
        "• Combined confidence score (0–100%)",
        "• Plain-English reason",
    ]
    for i, line in enumerate(left_lines):
        c = TXT3 if i == 0 else TXT2
        txt(slide, line, lx + Inches(0.2), Inches(1.68) + Inches(0.42 * i),
            Inches(5.3), Inches(0.38), sz=13, clr=c)

    rx = Inches(6.88)
    card(slide, rx, Inches(1.1), Inches(5.7), Inches(4.0))
    rect(slide, rx, Inches(1.1), Inches(0.05), Inches(4.0), ACC)
    txt(slide, "Audit Trail", rx + Inches(0.2), Inches(1.2),
        Inches(5.4), Inches(0.42), sz=16, bold=True, clr=TXT1)
    right_lines = [
        "Every action logged with:",
        "• UTC timestamp",
        "• Action type: criteria_extracted / bidder_processed",
        "  criterion_evaluated / human_review_action",
        "  vision_ocr_invoked / precomputed_fallback_used",
        "• Model version & Actor (system / officer)",
        "• Full payload JSON — Exportable as CSV",
    ]
    for i, line in enumerate(right_lines):
        c = TXT3 if i == 0 else TXT2
        txt(slide, line, rx + Inches(0.2), Inches(1.68) + Inches(0.42 * i),
            Inches(5.3), Inches(0.38), sz=13, clr=c)

    rect(slide, Inches(0.75), Inches(5.35), Inches(11.83), Inches(1.25),
         AMBBG, AMB, Pt(2))
    txt(slide, "The Safety Rule",
        Inches(0.95), Inches(5.42), Inches(5), Inches(0.38),
        sz=15, bold=True, clr=AMB)
    txt(slide, ("If combined confidence is 0.55–0.80 AND verdict is not_eligible, "
                "the verdict is automatically downgraded to needs_review. "
                "A bidder is NEVER silently disqualified at medium confidence."),
        Inches(0.95), Inches(5.85), Inches(11.5), Inches(0.65),
        sz=13, clr=TXT2)


# ── Slide 7 ─────────────────────────────────────────────────────────────────
def s7(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Demo: Three Bidders, Three Outcomes")
    bidders = [
        ("Bidder A — Apex Constructions Pvt. Ltd.", GRN, GRNBG, "ELIGIBLE",
         ["C1 Turnover: INR 6.37 Cr avg (threshold 5 Cr) — PASS",
          "C2 Projects: 5 completed incl. CRPF barracks — PASS",
          "C3 GST: GSTIN 27AABCA1234F1Z5, Active — PASS",
          "C4 ISO 9001:2015: Valid June 2027 — PASS",
          "All typed PDFs, confidence ≥ 93% on all criteria"]),
        ("Bidder B — BuildRight Enterprises", RED_C, REDBG, "NOT ELIGIBLE",
         ["C1 Turnover: INR 1.5 Cr avg (threshold 5 Cr) — FAIL",
          "\"INR 1.5 Cr is below required minimum of INR 5 Cr\"",
          "C2–C4: All pass",
          "Auto-disqualified at high confidence (95%)"]),
        ("Bidder C — Shree Constructions & Services", AMB, AMBBG, "NEEDS REVIEW",
         ["C1 Turnover: Submitted as blurry scan",
          "Tesseract ~55% → Vision LLM: INR 5.4 Cr",
          "Combined confidence 0.58 → safety rule triggered",
          "C2: Exactly 3 projects (borderline)",
          "C3–C4: Pass"]),
    ]
    cw, ch = Inches(3.85), Inches(3.6)
    for i, (name, color, cbg, verdict, bullets) in enumerate(bidders):
        sx = Inches(0.75) + i * Inches(4.2)
        card(slide, sx, Inches(1.1), cw, ch)
        rect(slide, sx, Inches(1.1), cw, Inches(0.07), color)
        txt(slide, name, sx + Inches(0.15), Inches(1.2), cw - Inches(0.3), Inches(0.5),
            sz=12, bold=True, clr=TXT1)
        chip(slide, sx + Inches(0.15), Inches(1.78), Inches(3.3), Inches(0.35),
             cbg, color, verdict)
        for j, b in enumerate(bullets):
            txt(slide, b, sx + Inches(0.15), Inches(2.24) + Inches(0.37 * j),
                cw - Inches(0.3), Inches(0.33), sz=11, clr=TXT2)

    # Metric strip
    metrics = [
        ("Criteria extracted", "5"),
        ("Bidder docs processed", "15"),
        ("LLM evaluation calls", "15"),
        ("Vision OCR invocations", "1"),
        ("Human review items", "1"),
        ("Total audit entries", "20+"),
    ]
    rect(slide, Inches(0.75), Inches(5.0), Inches(11.83), Inches(1.0),
         ACBG, ACC, Pt(1))
    mw = Inches(1.95)
    for i, (label, val) in enumerate(metrics):
        mx = Inches(0.85) + i * Inches(1.97)
        txt(slide, val, mx, Inches(5.08), mw, Inches(0.4),
            sz=22, bold=True, clr=ACC, align=PP_ALIGN.CENTER)
        txt(slide, label, mx, Inches(5.53), mw, Inches(0.35),
            sz=11, clr=TXT3, align=PP_ALIGN.CENTER)


# ── Slide 8 ─────────────────────────────────────────────────────────────────
def s8(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Stack, Impact & What's Next")
    stack = [
        ("UI & orchestration", "Streamlit 1.39"),
        ("LLM", "DeepSeek API (OpenAI-compatible)"),
        ("OCR Tier 1", "PyMuPDF 1.24"),
        ("OCR Tier 2", "Tesseract"),
        ("OCR Tier 3", "DeepSeek Vision LLM"),
        ("Semantic retrieval", "sentence-transformers all-MiniLM-L6-v2"),
        ("Data validation", "Pydantic v2"),
        ("Audit log", "SQLite"),
        ("Deployment", "Streamlit Cloud / HuggingFace Spaces"),
    ]
    rect(slide, Inches(0.75), Inches(1.05), Inches(6.0), Inches(0.4),
         ACBG, ACC, Pt(1))
    txt(slide, "Component", Inches(0.85), Inches(1.1), Inches(2.7), Inches(0.3),
        sz=13, bold=True, clr=ACC)
    txt(slide, "Technology", Inches(3.65), Inches(1.1), Inches(2.9), Inches(0.3),
        sz=13, bold=True, clr=ACC)
    for i, (comp, tech) in enumerate(stack):
        ry = Inches(1.45) + Inches(0.38 * i)
        bg_r = RGBColor(0xFA, 0xFB, 0xFF) if i % 2 == 0 else BG
        rect(slide, Inches(0.75), ry, Inches(6.0), Inches(0.38), bg_r)
        txt(slide, comp, Inches(0.85), ry + Inches(0.06), Inches(2.7), Inches(0.3),
            sz=12, clr=TXT3)
        txt(slide, tech, Inches(3.65), ry + Inches(0.06), Inches(2.9), Inches(0.3),
            sz=12, clr=TXT2)

    txt(slide, "What's Next", Inches(7.2), Inches(1.05), Inches(5.7), Inches(0.42),
        sz=16, bold=True, clr=TXT1)
    future = [
        "Multi-tender workspace — same bidder pool, multiple tenders",
        "GeM portal API integration — live tender ingestion",
        "Automated bidder ranking with weighted scoring",
        "LayoutLM for complex financial tables in scanned statements",
        "Multi-evaluator workflow with role-based approval",
        "Review queue email/SMS notifications",
        "Audit PDF export for procurement oversight submissions",
    ]
    for i, f in enumerate(future):
        txt(slide, f"• {f}", Inches(7.2), Inches(1.55) + Inches(0.47 * i),
            Inches(5.8), Inches(0.42), sz=13, clr=TXT2)

    rect(slide, Inches(0.75), Inches(6.18), Inches(11.83), Inches(0.95),
         ACBG, ACC, Pt(2))
    txt(slide, ("3–5 days → minutes.  Every verdict traceable to a document, page, and model version."
                "  Built in one hackathon session. Deployable today."),
        Inches(0.95), Inches(6.3), Inches(11.5), Inches(0.75),
        sz=14, bold=True, clr=ACC, align=PP_ALIGN.CENTER)


def main():
    os.makedirs("deck", exist_ok=True)
    prs = new_prs()
    s1(prs); s2(prs); s3(prs); s4(prs)
    s5(prs); s6(prs); s7(prs); s8(prs)
    out = "deck/TenderIQ_v2_clean_minimal.pptx"
    prs.save(out)
    print(f"OK: Saved {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
