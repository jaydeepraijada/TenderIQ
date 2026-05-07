#!/usr/bin/env python3
"""TenderIQ v3 — Government Official (PPTX)"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Palette
BLUE  = RGBColor(0x00, 0x35, 0x80)   # deep government blue
SAF   = RGBColor(0xFF, 0x99, 0x33)   # saffron
IGRN  = RGBColor(0x13, 0x88, 0x08)   # India green
BGOFF = RGBColor(0xF5, 0xF5, 0xF0)   # off-white
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TXT   = RGBColor(0x1A, 0x1A, 0x1A)
LTBLUE = RGBColor(0xEB, 0xF0, 0xFF)  # light blue fill
FONT_H = "Times New Roman"
FONT_B = "Arial"


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h, sz=14, bold=False, clr=None, align=PP_ALIGN.LEFT,
        font=None, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font or FONT_B
    if clr:
        r.font.color.rgb = clr
    return tb


def header_bar(slide, title):
    """Blue header band with title, saffron underline, footer."""
    rect(slide, 0, 0, Inches(13.33), Inches(0.9), BLUE)
    rect(slide, 0, Inches(0.9), Inches(13.33), Inches(0.06), SAF)
    txt(slide, title, Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.7),
        sz=22, bold=True, clr=WHITE, font=FONT_H)
    footer(slide)


def footer(slide):
    rect(slide, 0, Inches(7.32), Inches(13.33), Inches(0.18), BLUE)
    txt(slide, "TenderIQ  ·  CRPF Hackathon  ·  Theme 3",
        Inches(0.3), Inches(7.33), Inches(12.5), Inches(0.16),
        sz=9, clr=WHITE, align=PP_ALIGN.CENTER)


def bg(slide):
    rect(slide, 0, 0, Inches(13.33), Inches(7.5), BGOFF)


def callout(slide, x, y, w, h, text):
    rect(slide, x, y, w, h, LTBLUE, BLUE, Pt(1))
    txt(slide, text, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), h - Inches(0.2),
        sz=13, clr=TXT)


# ── Slide 1 ─────────────────────────────────────────────────────────────────
def s1(prs):
    slide = blank(prs)
    bg(slide)
    rect(slide, 0, 0, Inches(13.33), Inches(1.2), BLUE)
    rect(slide, 0, Inches(1.2), Inches(13.33), Inches(0.08), SAF)
    txt(slide, "Government of India  ·  Ministry of Home Affairs",
        Inches(0.5), Inches(0.1), Inches(12.33), Inches(0.4),
        sz=13, clr=WHITE, align=PP_ALIGN.CENTER, font=FONT_H)
    txt(slide, "Central Reserve Police Force",
        Inches(0.5), Inches(0.55), Inches(12.33), Inches(0.5),
        sz=18, bold=True, clr=WHITE, align=PP_ALIGN.CENTER, font=FONT_H)
    rect(slide, Inches(0), Inches(1.28), Inches(13.33), Inches(0.04), SAF)
    txt(slide, "TenderIQ",
        Inches(0.5), Inches(1.8), Inches(12.33), Inches(1.4),
        sz=64, bold=True, clr=BLUE, align=PP_ALIGN.CENTER, font=FONT_H)
    txt(slide, "Explainable AI for Government Tender Evaluation",
        Inches(0.5), Inches(3.3), Inches(12.33), Inches(0.65),
        sz=22, bold=True, clr=TXT, align=PP_ALIGN.CENTER, font=FONT_H)
    rect(slide, Inches(3), Inches(4.1), Inches(7.33), Inches(0.04), SAF)
    txt(slide, "CRPF Hackathon  —  Theme 3: AI-Based Tender Evaluation",
        Inches(0.5), Inches(4.3), Inches(12.33), Inches(0.5),
        sz=16, clr=BLUE, align=PP_ALIGN.CENTER, font=FONT_H)
    txt(slide, "From days to minutes. Every decision traceable.",
        Inches(0.5), Inches(5.0), Inches(12.33), Inches(0.45),
        sz=14, italic=True, clr=TXT, align=PP_ALIGN.CENTER)
    footer(slide)


# ── Slide 2 ─────────────────────────────────────────────────────────────────
def s2(prs):
    slide = blank(prs)
    bg(slide)
    header_bar(slide, "The Problem with Manual Tender Evaluation")
    callouts_data = [
        ("3–5 Days", "per tender evaluation by committee"),
        ("Inconsistent", "two evaluators, two different conclusions"),
        ("No Audit Trail", "decisions made untraceably"),
    ]
    cw, ch = Inches(3.7), Inches(2.0)
    for i, (big, sub) in enumerate(callouts_data):
        sx = Inches(0.6) + i * Inches(4.1)
        rect(slide, sx, Inches(1.2), cw, ch, WHITE, BLUE, Pt(2))
        rect(slide, sx, Inches(1.2), cw, Inches(0.08), SAF)
        txt(slide, big, sx + Inches(0.15), Inches(1.4), cw - Inches(0.3), Inches(0.8),
            sz=28, bold=True, clr=BLUE, align=PP_ALIGN.CENTER, font=FONT_H)
        txt(slide, sub, sx + Inches(0.15), Inches(2.25), cw - Inches(0.3), Inches(0.7),
            sz=13, clr=TXT, align=PP_ALIGN.CENTER)
    rect(slide, Inches(0.6), Inches(3.45), Inches(12.13), Inches(0.04), SAF)
    bullets = [
        "• Mixed document formats: typed PDFs, scanned certificates, phone photographs",
        "• Government procurement worth ₹50 lakh crore+ annually in India",
        "• Project execution delays traced directly to procurement bottlenecks",
    ]
    for i, b in enumerate(bullets):
        txt(slide, b, Inches(0.75), Inches(3.65) + Inches(0.45 * i),
            Inches(12), Inches(0.38), sz=14, clr=TXT)


# ── Slide 3 ─────────────────────────────────────────────────────────────────
def s3(prs):
    slide = blank(prs)
    bg(slide)
    header_bar(slide, "TenderIQ — Four Stages, End to End")
    stages = [
        ("Stage 1", "Extract",
         "DeepSeek LLM reads the tender PDF\nand returns each criterion as structured\nJSON: category, mandatory flag, threshold\nrule, source clause, query hints."),
        ("Stage 2", "OCR & Index",
         "Three-tier pipeline handles any document\nformat. All text chunked and indexed for\nsemantic retrieval."),
        ("Stage 3", "Evaluate",
         "Vector search finds relevant evidence.\nLLM produces a verdict with confidence.\nSafety rule prevents silent disqualification."),
        ("Stage 4", "Review & Audit",
         "Borderline cases go to human review\nqueue. Every action logged. Full audit\ntrail exportable as CSV."),
    ]
    cw, ch = Inches(2.9), Inches(3.8)
    for i, (num, title, body) in enumerate(stages):
        sx = Inches(0.6) + i * Inches(3.15)
        rect(slide, sx, Inches(1.1), cw, ch, WHITE, BLUE, Pt(1))
        rect(slide, sx, Inches(1.1), cw, Inches(0.07), BLUE)
        txt(slide, num, sx + Inches(0.15), Inches(1.2), cw - Inches(0.3), Inches(0.38),
            sz=11, bold=True, clr=SAF, font=FONT_H)
        txt(slide, title, sx + Inches(0.15), Inches(1.6), cw - Inches(0.3), Inches(0.45),
            sz=16, bold=True, clr=BLUE, font=FONT_H)
        txt(slide, body, sx + Inches(0.15), Inches(2.1), cw - Inches(0.3), Inches(2.5),
            sz=12, clr=TXT)
    rect(slide, Inches(0.6), Inches(5.1), Inches(12.13), Inches(0.85),
         LTBLUE, BLUE, Pt(1))
    txt(slide, '"Minutes, not days. Every verdict traceable to a document and page."',
        Inches(0.8), Inches(5.22), Inches(11.9), Inches(0.65),
        sz=14, bold=True, clr=BLUE, align=PP_ALIGN.CENTER, font=FONT_H)


# ── Slide 4 ─────────────────────────────────────────────────────────────────
def s4(prs):
    slide = blank(prs)
    bg(slide)
    header_bar(slide, "System Architecture")

    bw, bh = Inches(2.5), Inches(0.52)

    def abox(x, y, label):
        rect(slide, x, y, bw, bh, WHITE, BLUE, Pt(1))
        txt(slide, label, x + Inches(0.1), y + Inches(0.06),
            bw - Inches(0.2), bh - Inches(0.1),
            sz=11, clr=TXT, align=PP_ALIGN.CENTER)

    def arr(x, y):
        rect(slide, x + Inches(1.15), y, Inches(0.2), Inches(0.28), BLUE)

    lx = Inches(0.6)
    abox(lx, Inches(1.05), "Tender PDF")
    arr(lx, Inches(1.57))
    abox(lx, Inches(1.85), "DeepSeek LLM\n(Extract Criteria)")
    arr(lx, Inches(2.37))
    abox(lx, Inches(2.65), "Criteria JSON\n(C1–C5 structured)")

    rx = Inches(4.0)
    abox(rx, Inches(1.05), "Bidder Docs\n(PDFs · scans · photos)")
    arr(rx, Inches(1.57))
    abox(rx, Inches(1.85), "3-Tier OCR Pipeline\n① PyMuPDF ② Tesseract ③ Vision LLM")
    arr(rx, Inches(2.37))
    abox(rx, Inches(2.65), "Vector Index\n(all-MiniLM-L6-v2 embeddings)")

    cx = Inches(3.15)
    arr(cx, Inches(3.17))
    abox(Inches(2.1), Inches(3.45), "DeepSeek LLM — Evaluate each criterion")

    rect(slide, Inches(0.9), Inches(4.2), Inches(2.0), Inches(0.52),
         RGBColor(0xD1, 0xFA, 0xE5), IGRN, Pt(1))
    txt(slide, "ELIGIBLE /\nNOT ELIGIBLE", Inches(0.9), Inches(4.2), Inches(2.0), Inches(0.52),
        sz=11, clr=IGRN, align=PP_ALIGN.CENTER, bold=True)

    rect(slide, Inches(3.5), Inches(4.2), Inches(2.0), Inches(0.52),
         RGBColor(0xFE, 0xF3, 0xC7), SAF, Pt(1))
    txt(slide, "UNDER REVIEW\nHuman Review Queue", Inches(3.5), Inches(4.2), Inches(2.0), Inches(0.52),
        sz=11, clr=RGBColor(0xB4, 0x5A, 0x09), align=PP_ALIGN.CENTER, bold=True)

    abox(Inches(2.1), Inches(5.0), "SQLite Audit Log")

    # Key facts
    rect(slide, Inches(7.2), Inches(1.05), Inches(5.7), Inches(5.2), WHITE, BLUE, Pt(1))
    rect(slide, Inches(7.2), Inches(1.05), Inches(5.7), Inches(0.06), BLUE)
    txt(slide, "Key Technical Facts", Inches(7.35), Inches(1.15), Inches(5.4), Inches(0.42),
        sz=14, bold=True, clr=BLUE, font=FONT_H)
    facts = [
        "Single-process Streamlit app — no separate backend",
        "Deployable to Streamlit Cloud or HuggingFace Spaces",
        "All storage is local: SQLite + in-memory vector index",
        "Only external dependency: DeepSeek API",
    ]
    for i, f in enumerate(facts):
        txt(slide, f"• {f}", Inches(7.35), Inches(1.72) + Inches(0.65 * i),
            Inches(5.4), Inches(0.55), sz=13, clr=TXT)


# ── Slide 5 ─────────────────────────────────────────────────────────────────
def s5(prs):
    slide = blank(prs)
    bg(slide)
    header_bar(slide, "Three-Tier OCR — Handling Any Document Format")
    tiers = [
        ("Tier 1 — PyMuPDF",
         "Trigger: Typed / digital PDF\nCost: Free, instant\nConfidence: 1.0 (lossless)\nSource label: Typed PDF",
         BLUE),
        ("Tier 2 — Tesseract",
         "Trigger: Scanned PDF or image\nCost: Free, local, fast\nConfidence: Mean per-word OCR\nSource label: Tesseract",
         RGBColor(0x7C, 0x3A, 0xED)),
        ("Tier 3 — DeepSeek Vision LLM",
         "Trigger: Tesseract confidence < 65%\nCost: One API call\nConfidence: 0.95\nSource label: Vision LLM\nLogged: vision_ocr_invoked",
         SAF),
    ]
    cw, ch = Inches(3.6), Inches(3.1)
    for i, (title, body, accent) in enumerate(tiers):
        sx = Inches(0.6) + i * Inches(4.25)
        rect(slide, sx, Inches(1.1), cw, ch, WHITE, BLUE, Pt(1))
        rect(slide, sx, Inches(1.1), cw, Inches(0.08), accent)
        txt(slide, title, sx + Inches(0.15), Inches(1.25), cw - Inches(0.3), Inches(0.48),
            sz=14, bold=True, clr=accent, font=FONT_H)
        txt(slide, body, sx + Inches(0.15), Inches(1.78), cw - Inches(0.3), Inches(2.2),
            sz=13, clr=TXT)

    callout(slide, Inches(0.6), Inches(4.5), Inches(12.13), Inches(1.85),
            "Demo Scenario — Bidder C submits a blurry, rotated CA certificate scan.\n"
            "Tesseract reads it at ~55% confidence. Vision LLM transcribes the turnover figure correctly.\n"
            "Combined confidence = 0.58 → routed to human review.\n"
            "This is intentional — borderline evidence requires a human.")


# ── Slide 6 ─────────────────────────────────────────────────────────────────
def s6(prs):
    slide = blank(prs)
    bg(slide)
    header_bar(slide, "Every Decision is Explainable and Auditable")

    lx = Inches(0.6)
    rect(slide, lx, Inches(1.1), Inches(5.85), Inches(4.0), WHITE, BLUE, Pt(1))
    rect(slide, lx, Inches(1.1), Inches(5.85), Inches(0.07), BLUE)
    txt(slide, "Criterion-Level Verdicts", lx + Inches(0.15), Inches(1.2),
        Inches(5.5), Inches(0.42), sz=14, bold=True, clr=BLUE, font=FONT_H)
    for i, line in enumerate([
        "Each (bidder × criterion) pair shows:",
        "• Which criterion was checked",
        "• Which document and page provided evidence",
        "• What value was extracted (e.g. 'INR 6.2 Cr')",
        "• Which OCR tier read the document",
        "• Combined confidence score (0–100%)",
        "• Plain-English reason",
    ]):
        txt(slide, line, lx + Inches(0.15), Inches(1.65) + Inches(0.42 * i),
            Inches(5.55), Inches(0.38), sz=13, clr=TXT)

    rx = Inches(6.98)
    rect(slide, rx, Inches(1.1), Inches(5.85), Inches(4.0), WHITE, BLUE, Pt(1))
    rect(slide, rx, Inches(1.1), Inches(5.85), Inches(0.07), BLUE)
    txt(slide, "Audit Trail", rx + Inches(0.15), Inches(1.2),
        Inches(5.5), Inches(0.42), sz=14, bold=True, clr=BLUE, font=FONT_H)
    for i, line in enumerate([
        "Every action logged with:",
        "• UTC timestamp",
        "• Action type: criteria_extracted / bidder_processed",
        "  criterion_evaluated / human_review_action",
        "  vision_ocr_invoked / precomputed_fallback_used",
        "• Model version & Actor (system / officer)",
        "• Full payload JSON — Exportable as CSV",
    ]):
        txt(slide, line, rx + Inches(0.15), Inches(1.65) + Inches(0.42 * i),
            Inches(5.55), Inches(0.38), sz=13, clr=TXT)

    rect(slide, Inches(0.6), Inches(5.35), Inches(12.13), Inches(1.2),
         LTBLUE, BLUE, Pt(2))
    txt(slide, "The Safety Rule:",
        Inches(0.8), Inches(5.42), Inches(5), Inches(0.38),
        sz=14, bold=True, clr=BLUE, font=FONT_H)
    txt(slide, ("If combined confidence is 0.55–0.80 AND verdict is NOT ELIGIBLE, "
                "the verdict is automatically downgraded to UNDER REVIEW. "
                "A bidder is NEVER silently disqualified at medium confidence."),
        Inches(0.8), Inches(5.85), Inches(11.9), Inches(0.65),
        sz=13, clr=TXT)


# ── Slide 7 ─────────────────────────────────────────────────────────────────
def s7(prs):
    slide = blank(prs)
    bg(slide)
    header_bar(slide, "Demo: Three Bidders, Three Outcomes")
    bidders = [
        ("Bidder A — Apex Constructions Pvt. Ltd.", IGRN, "ELIGIBLE",
         ["C1 Turnover: INR 6.37 Cr avg (threshold 5 Cr) — PASS",
          "C2 Projects: 5 completed incl. CRPF barracks — PASS",
          "C3 GST: GSTIN 27AABCA1234F1Z5, Active — PASS",
          "C4 ISO 9001:2015: Valid June 2027 — PASS",
          "All typed PDFs, confidence ≥ 93%"]),
        ("Bidder B — BuildRight Enterprises", RGBColor(0xB9, 0x1C, 0x1C), "NOT ELIGIBLE",
         ["C1 Turnover: INR 1.5 Cr avg (threshold 5 Cr) — FAIL",
          "Reason: INR 1.5 Cr is below required INR 5 Cr",
          "C2–C4: All pass",
          "Auto-disqualified at high confidence (95%)"]),
        ("Bidder C — Shree Constructions & Services", SAF, "UNDER REVIEW",
         ["C1 Turnover: Submitted as blurry scan",
          "Tesseract ~55% → Vision LLM: INR 5.4 Cr",
          "Combined confidence 0.58 → safety rule triggered",
          "C2: Exactly 3 projects (borderline)",
          "C3–C4: Pass"]),
    ]
    cw, ch = Inches(3.85), Inches(3.5)
    for i, (name, color, verdict, bullets) in enumerate(bidders):
        sx = Inches(0.6) + i * Inches(4.25)
        rect(slide, sx, Inches(1.1), cw, ch, WHITE, BLUE, Pt(1))
        rect(slide, sx, Inches(1.1), cw, Inches(0.08), color)
        txt(slide, name, sx + Inches(0.15), Inches(1.22), cw - Inches(0.3), Inches(0.5),
            sz=12, bold=True, clr=BLUE, font=FONT_H)
        rect(slide, sx + Inches(0.15), Inches(1.8), Inches(3.4), Inches(0.35),
             RGBColor(0xEB, 0xF0, 0xFF), color, Pt(1))
        txt(slide, verdict, sx + Inches(0.15), Inches(1.8), Inches(3.4), Inches(0.35),
            sz=12, bold=True, clr=color, align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            txt(slide, b, sx + Inches(0.15), Inches(2.25) + Inches(0.37 * j),
                cw - Inches(0.3), Inches(0.33), sz=11, clr=TXT)

    # Metric strip
    metrics = [
        ("Criteria extracted", "5"),
        ("Bidder docs processed", "15"),
        ("LLM evaluation calls", "15"),
        ("Vision OCR invocations", "1"),
        ("Human review items", "1"),
        ("Total audit entries", "20+"),
    ]
    rect(slide, Inches(0.6), Inches(4.9), Inches(12.13), Inches(1.0),
         LTBLUE, BLUE, Pt(1))
    mw = Inches(2.0)
    for i, (label, val) in enumerate(metrics):
        mx = Inches(0.75) + i * Inches(2.0)
        txt(slide, val, mx, Inches(4.97), mw, Inches(0.38),
            sz=22, bold=True, clr=BLUE, align=PP_ALIGN.CENTER, font=FONT_H)
        txt(slide, label, mx, Inches(5.42), mw, Inches(0.35),
            sz=11, clr=TXT, align=PP_ALIGN.CENTER)


# ── Slide 8 ─────────────────────────────────────────────────────────────────
def s8(prs):
    slide = blank(prs)
    bg(slide)
    header_bar(slide, "Stack, Impact & What's Next")
    stack = [
        ("UI & orchestration", "Streamlit 1.39"),
        ("LLM", "DeepSeek API (OpenAI-compatible)"),
        ("OCR Tier 1", "PyMuPDF 1.24"),
        ("OCR Tier 2", "Tesseract"),
        ("OCR Tier 3", "DeepSeek Vision LLM"),
        ("Semantic retrieval", "sentence-transformers all-MiniLM-L6-v2"),
        ("Data validation", "Pydantic v2"),
        ("Audit log", "SQLite"),
        ("Deployment", "Streamlit Cloud / HuggingFace Spaces"),
    ]
    rect(slide, Inches(0.6), Inches(1.05), Inches(6.0), Inches(0.4),
         BLUE)
    txt(slide, "Component", Inches(0.7), Inches(1.1), Inches(2.7), Inches(0.3),
        sz=13, bold=True, clr=WHITE, font=FONT_H)
    txt(slide, "Technology", Inches(3.55), Inches(1.1), Inches(2.9), Inches(0.3),
        sz=13, bold=True, clr=WHITE, font=FONT_H)
    for i, (comp, tech) in enumerate(stack):
        ry = Inches(1.45) + Inches(0.38 * i)
        bg_r = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xFF)
        rect(slide, Inches(0.6), ry, Inches(6.0), Inches(0.38), bg_r,
             RGBColor(0xCB, 0xD5, 0xE1), Pt(0))
        txt(slide, comp, Inches(0.7), ry + Inches(0.06), Inches(2.7), Inches(0.3),
            sz=12, clr=TXT)
        txt(slide, tech, Inches(3.55), ry + Inches(0.06), Inches(2.9), Inches(0.3),
            sz=12, clr=TXT)

    txt(slide, "Future Work", Inches(7.2), Inches(1.05), Inches(5.7), Inches(0.42),
        sz=15, bold=True, clr=BLUE, font=FONT_H)
    future = [
        "Multi-tender workspace — same bidder pool, multiple tenders",
        "GeM portal API integration — live tender ingestion",
        "Automated bidder ranking with weighted scoring",
        "LayoutLM for complex financial tables in scanned statements",
        "Multi-evaluator workflow with role-based approval",
        "Review queue email/SMS notifications",
        "Audit PDF export for procurement oversight submissions",
    ]
    for i, f in enumerate(future):
        txt(slide, f"• {f}", Inches(7.2), Inches(1.55) + Inches(0.47 * i),
            Inches(5.8), Inches(0.42), sz=13, clr=TXT)

    rect(slide, Inches(0.6), Inches(6.18), Inches(12.13), Inches(0.95),
         LTBLUE, BLUE, Pt(2))
    txt(slide, ("3–5 days → minutes.  Every verdict traceable to a document, page, and model version."
                "  Built in one hackathon session. Deployable today."),
        Inches(0.8), Inches(6.3), Inches(11.9), Inches(0.75),
        sz=13, bold=True, clr=BLUE, align=PP_ALIGN.CENTER, font=FONT_H)


def main():
    os.makedirs("deck", exist_ok=True)
    prs = new_prs()
    s1(prs); s2(prs); s3(prs); s4(prs)
    s5(prs); s6(prs); s7(prs); s8(prs)
    out = "deck/TenderIQ_v3_government_official.pptx"
    prs.save(out)
    print(f"OK: Saved {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
