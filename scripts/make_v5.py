#!/usr/bin/env python3
"""TenderIQ v5 — Data Forward (PPTX with embedded matplotlib charts)"""
import os
import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Palette
BG    = RGBColor(0xFA, 0xFA, 0xFA)
PRI   = RGBColor(0x1E, 0x29, 0x3B)
ACC   = RGBColor(0x63, 0x66, 0xF1)   # indigo
TXT   = RGBColor(0x33, 0x41, 0x55)
TXT2  = RGBColor(0x64, 0x74, 0x8B)
GRD   = RGBColor(0xE2, 0xE8, 0xF0)
GRN   = RGBColor(0x22, 0xC5, 0x5E)
RED_C = RGBColor(0xEF, 0x44, 0x44)
AMB   = RGBColor(0xF5, 0x9E, 0x0B)
BLU   = RGBColor(0x3B, 0x82, 0xF6)
PUR   = RGBColor(0x8B, 0x5C, 0xF6)
FONT  = "Calibri"

HEX = {
    "bg":   "#FAFAFA",
    "pri":  "#1E293B",
    "acc":  "#6366F1",
    "txt":  "#334155",
    "txt2": "#64748B",
    "grd":  "#E2E8F0",
    "grn":  "#22C55E",
    "red":  "#EF4444",
    "amb":  "#F59E0B",
    "blu":  "#3B82F6",
    "pur":  "#8B5CF6",
}


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, text, x, y, w, h, sz=14, bold=False, clr=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.name = FONT
    if clr:
        r.font.color.rgb = clr
    return tb


def bg(slide):
    rect(slide, 0, 0, Inches(13.33), Inches(7.5), BG)


def heading(slide, text, y=Inches(0.3)):
    rect(slide, Inches(0.6), y + Inches(0.08), Inches(0.06), Inches(0.44), ACC)
    txt(slide, text, Inches(0.8), y, Inches(12.3), Inches(0.6),
        sz=24, bold=True, clr=ACC)


def png_to_slide(slide, fig, x, y, w, h):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=HEX["bg"], transparent=True)
    plt.close(fig)
    buf.seek(0)
    slide.shapes.add_picture(buf, x, y, w, h)


def mpl_style():
    plt.rcParams.update({
        "axes.facecolor": HEX["bg"],
        "figure.facecolor": HEX["bg"],
        "axes.edgecolor": HEX["grd"],
        "axes.grid": True,
        "grid.color": HEX["grd"],
        "grid.linewidth": 0.7,
        "text.color": HEX["pri"],
        "xtick.color": HEX["txt2"],
        "ytick.color": HEX["txt2"],
        "font.family": "DejaVu Sans",
    })


# ── Slide 1 ─────────────────────────────────────────────────────────────────
def s1(prs):
    slide = blank(prs)
    bg(slide)
    rect(slide, 0, 0, Inches(0.5), Inches(7.5), ACC)
    txt(slide, "TenderIQ", Inches(0.7), Inches(1.8), Inches(11.5), Inches(1.5),
        sz=64, bold=True, clr=PRI, align=PP_ALIGN.LEFT)
    txt(slide, "Explainable AI for Government Tender Evaluation",
        Inches(0.7), Inches(3.3), Inches(11.5), Inches(0.65),
        sz=22, clr=ACC)
    rect(slide, Inches(0.7), Inches(4.1), Inches(8), Inches(0.05), GRD)
    txt(slide, "CRPF Hackathon  ·  Theme 3",
        Inches(0.7), Inches(4.3), Inches(8), Inches(0.5),
        sz=16, clr=TXT2)
    txt(slide, "From days to minutes. Every decision traceable.",
        Inches(0.7), Inches(4.95), Inches(9), Inches(0.5),
        sz=15, clr=TXT)


# ── Slide 2 ─────────────────────────────────────────────────────────────────
def s2(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "The Problem with Manual Tender Evaluation")
    mpl_style()
    fig, ax = plt.subplots(figsize=(5.5, 3.8))
    categories = ["Manual\nProcess", "TenderIQ"]
    values = [4.0, 0.02]
    colors_ = [HEX["red"], HEX["grn"]]
    bars = ax.bar(categories, values, color=colors_, width=0.45, zorder=3)
    ax.set_ylabel("Days per tender", color=HEX["txt2"], fontsize=11)
    ax.set_ylim(0, 5.5)
    ax.set_title("Evaluation time comparison", color=HEX["pri"], fontsize=12, pad=8)
    for bar, val in zip(bars, values):
        label = f"{val} days" if val >= 0.1 else "Minutes"
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.1,
                label, ha="center", va="bottom", fontsize=12,
                fontweight="bold", color=HEX["pri"])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(axis="x", labelsize=12)
    png_to_slide(slide, fig, Inches(0.6), Inches(1.0), Inches(5.8), Inches(4.5))

    bullets = [
        ("3–5 Days", "per tender evaluation by committee"),
        ("Inconsistent", "two evaluators, two conclusions"),
        ("No Audit Trail", "decisions made untraceably"),
    ]
    for i, (big, sub) in enumerate(bullets):
        rect(slide, Inches(7.3), Inches(1.2) + Inches(1.55 * i),
             Inches(5.5), Inches(1.3), RGBColor(0xF1, 0xF5, 0xF9), GRD, Pt(1))
        txt(slide, big, Inches(7.5), Inches(1.3) + Inches(1.55 * i),
            Inches(5.0), Inches(0.55), sz=22, bold=True, clr=ACC)
        txt(slide, sub, Inches(7.5), Inches(1.85) + Inches(1.55 * i),
            Inches(5.0), Inches(0.45), sz=13, clr=TXT)
    txt(slide, ("• Mixed document formats: typed PDFs, scans, phone photographs\n"
                "• Government procurement worth ₹50 lakh crore+ annually in India"),
        Inches(7.3), Inches(6.0), Inches(5.5), Inches(0.9), sz=12, clr=TXT2)


# ── Slide 3 ─────────────────────────────────────────────────────────────────
def s3(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "TenderIQ — Four Stages, End to End")
    stages = [
        ("1", "Extract",
         "DeepSeek LLM reads tender PDF → structured JSON criteria"),
        ("2", "OCR & Index",
         "3-tier pipeline handles any format → vector index"),
        ("3", "Evaluate",
         "Vector search + LLM verdict + safety rule"),
        ("4", "Review & Audit",
         "Human review queue + full exportable audit trail"),
    ]
    cw = Inches(2.8)
    for i, (num, title, body) in enumerate(stages):
        sx = Inches(0.6) + i * Inches(3.18)
        rect(slide, sx, Inches(1.1), cw, Inches(3.5), RGBColor(0xF1, 0xF5, 0xF9), GRD, Pt(1))
        # Number circle (drawn as square for simplicity)
        rect(slide, sx + Inches(0.15), Inches(1.25), Inches(0.65), Inches(0.65),
             ACC)
        txt(slide, num, sx + Inches(0.15), Inches(1.25), Inches(0.65), Inches(0.65),
            sz=22, bold=True, clr=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        txt(slide, title, sx + Inches(0.9), Inches(1.3), cw - Inches(1.05), Inches(0.5),
            sz=17, bold=True, clr=PRI)
        txt(slide, body, sx + Inches(0.15), Inches(2.0), cw - Inches(0.3), Inches(2.4),
            sz=13, clr=TXT)
        # Connector arrow (not last)
        if i < 3:
            txt(slide, "→", sx + cw + Inches(0.05), Inches(2.5), Inches(0.3), Inches(0.5),
                sz=22, bold=True, clr=ACC, align=PP_ALIGN.CENTER)

    rect(slide, Inches(0.6), Inches(5.0), Inches(12.13), Inches(0.85),
         RGBColor(0xEE, 0xF2, 0xFF), ACC, Pt(1))
    txt(slide, '"Minutes, not days. Every verdict traceable to a document and page."',
        Inches(0.8), Inches(5.12), Inches(11.9), Inches(0.65),
        sz=15, bold=True, clr=ACC, align=PP_ALIGN.CENTER)


# ── Slide 4 ─────────────────────────────────────────────────────────────────
def s4(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "System Architecture")

    bw, bh = Inches(2.5), Inches(0.52)

    def abox(x, y, label, clr=RGBColor(0xF1, 0xF5, 0xF9)):
        rect(slide, x, y, bw, bh, clr, GRD, Pt(1))
        txt(slide, label, x + Inches(0.1), y + Inches(0.06),
            bw - Inches(0.2), bh - Inches(0.1), sz=12, clr=TXT, align=PP_ALIGN.CENTER)

    def arr(x, y):
        rect(slide, x + Inches(1.15), y, Inches(0.2), Inches(0.28), ACC)

    lx = Inches(0.6)
    abox(lx, Inches(1.0), "Tender PDF")
    arr(lx, Inches(1.52))
    abox(lx, Inches(1.8), "DeepSeek LLM\n(Extract Criteria)")
    arr(lx, Inches(2.32))
    abox(lx, Inches(2.6), "Criteria JSON\n(C1–C5 structured)")

    rx = Inches(4.0)
    abox(rx, Inches(1.0), "Bidder Docs\n(PDFs · scans · photos)")
    arr(rx, Inches(1.52))
    abox(rx, Inches(1.8), "3-Tier OCR\n① PyMuPDF ② Tesseract ③ Vision LLM")
    arr(rx, Inches(2.32))
    abox(rx, Inches(2.6), "Vector Index\n(all-MiniLM-L6-v2)")

    cx = Inches(3.15)
    arr(cx, Inches(3.12))
    abox(Inches(2.1), Inches(3.4), "DeepSeek LLM — Evaluate each criterion",
         RGBColor(0xEE, 0xF2, 0xFF))

    rect(slide, Inches(0.9), Inches(4.15), Inches(2.0), Inches(0.52),
         RGBColor(0xDC, 0xFC, 0xE7), GRN, Pt(1))
    txt(slide, "eligible / not_eligible", Inches(0.9), Inches(4.15), Inches(2.0), Inches(0.52),
        sz=11, clr=GRN, align=PP_ALIGN.CENTER)

    rect(slide, Inches(3.5), Inches(4.15), Inches(2.2), Inches(0.52),
         RGBColor(0xFE, 0xF3, 0xC7), AMB, Pt(1))
    txt(slide, "needs_review\nHuman Review Queue", Inches(3.5), Inches(4.15), Inches(2.2), Inches(0.52),
        sz=11, clr=AMB, align=PP_ALIGN.CENTER)

    abox(Inches(2.1), Inches(5.0), "SQLite Audit Log")

    rect(slide, Inches(7.2), Inches(1.0), Inches(5.7), Inches(5.0),
         RGBColor(0xF8, 0xF9, 0xFF), ACC, Pt(1))
    txt(slide, "Key Technical Facts", Inches(7.35), Inches(1.1), Inches(5.4), Inches(0.4),
        sz=14, bold=True, clr=ACC)
    facts = [
        "Single-process Streamlit app — no separate backend",
        "Deployable to Streamlit Cloud or HuggingFace Spaces",
        "All storage is local: SQLite + in-memory vector index",
        "Only external dependency: DeepSeek API",
    ]
    for i, f in enumerate(facts):
        txt(slide, f"• {f}", Inches(7.35), Inches(1.65) + Inches(0.65 * i),
            Inches(5.4), Inches(0.55), sz=13, clr=TXT)


# ── Slide 5 ─────────────────────────────────────────────────────────────────
def s5(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Three-Tier OCR — Handling Any Document Format")
    mpl_style()

    # OCR confidence bar chart
    fig, ax = plt.subplots(figsize=(4.8, 3.2))
    tiers = ["Tier 1\nPyMuPDF", "Tier 2\nTesseract", "Tier 3\nVision LLM"]
    vals = [100, 60, 95]
    colors_ = [HEX["blu"], HEX["pur"], HEX["amb"]]
    bars = ax.bar(tiers, vals, color=colors_, width=0.5, zorder=3)
    ax.set_ylabel("Confidence (%)", color=HEX["txt2"], fontsize=10)
    ax.set_ylim(0, 110)
    ax.set_title("OCR confidence by tier", color=HEX["pri"], fontsize=11, pad=6)
    for bar, val in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1,
                f"{val}%", ha="center", va="bottom", fontsize=11,
                fontweight="bold", color=HEX["pri"])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    png_to_slide(slide, fig, Inches(0.6), Inches(1.0), Inches(5.0), Inches(3.8))

    tiers_info = [
        ("Tier 1 — PyMuPDF", "Trigger: Typed / digital PDF\nCost: Free, instant\nConf: 1.0 (lossless)\nLabel: Typed PDF",
         RGBColor(0x3B, 0x82, 0xF6)),
        ("Tier 2 — Tesseract", "Trigger: Scanned PDF / image\nCost: Free, local\nConf: Mean per-word scores\nLabel: Tesseract",
         RGBColor(0x8B, 0x5C, 0xF6)),
        ("Tier 3 — Vision LLM", "Trigger: Tesseract conf < 65%\nCost: One API call\nConf: 0.95\nLabel: Vision LLM",
         AMB),
    ]
    cw = Inches(2.55)
    for i, (title, body, accent) in enumerate(tiers_info):
        sx = Inches(6.1) + i * Inches(2.35)
        rect(slide, sx, Inches(1.0), cw, Inches(3.5), RGBColor(0xF8, 0xF9, 0xFF), accent, Pt(2))
        rect(slide, sx, Inches(1.0), cw, Inches(0.08), accent)
        txt(slide, title, sx + Inches(0.12), Inches(1.12), cw - Inches(0.24), Inches(0.45),
            sz=13, bold=True, clr=accent)
        txt(slide, body, sx + Inches(0.12), Inches(1.62), cw - Inches(0.24), Inches(2.5),
            sz=12, clr=TXT)

    rect(slide, Inches(0.6), Inches(5.0), Inches(12.13), Inches(1.75),
         RGBColor(0xEE, 0xF2, 0xFF), ACC, Pt(1))
    txt(slide, "Demo Scenario",
        Inches(0.8), Inches(5.1), Inches(8), Inches(0.35),
        sz=13, bold=True, clr=ACC)
    demo = ("Bidder C submits a blurry, rotated CA certificate scan. "
            "Tesseract reads it at ~55% confidence. Vision LLM transcribes the turnover figure correctly. "
            "Combined confidence = 0.58 → routed to human review. "
            "Borderline evidence requires a human — this is intentional.")
    txt(slide, demo, Inches(0.8), Inches(5.52), Inches(12.0), Inches(1.1), sz=13, clr=TXT)


# ── Slide 6 ─────────────────────────────────────────────────────────────────
def s6(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Every Decision is Explainable and Auditable")

    lx = Inches(0.6)
    rect(slide, lx, Inches(1.1), Inches(5.85), Inches(4.0),
         RGBColor(0xF8, 0xF9, 0xFF), ACC, Pt(1))
    rect(slide, lx, Inches(1.1), Inches(0.06), Inches(4.0), ACC)
    txt(slide, "Criterion-Level Verdicts", lx + Inches(0.2), Inches(1.2),
        Inches(5.5), Inches(0.42), sz=15, bold=True, clr=ACC)
    left_lines = [
        "Each (bidder × criterion) pair shows:",
        "• Which criterion was checked",
        "• Which document and page provided evidence",
        "• What value was extracted (e.g. 'INR 6.2 Cr')",
        "• Which OCR tier read the document",
        "• Combined confidence score (0–100%)",
        "• Plain-English reason",
    ]
    for i, line in enumerate(left_lines):
        c2 = TXT2 if i == 0 else TXT
        txt(slide, line, lx + Inches(0.2), Inches(1.68) + Inches(0.42 * i),
            Inches(5.5), Inches(0.38), sz=13, clr=c2)

    rx = Inches(6.98)
    rect(slide, rx, Inches(1.1), Inches(5.85), Inches(4.0),
         RGBColor(0xF8, 0xF9, 0xFF), ACC, Pt(1))
    rect(slide, rx, Inches(1.1), Inches(0.06), Inches(4.0), ACC)
    txt(slide, "Audit Trail", rx + Inches(0.2), Inches(1.2),
        Inches(5.5), Inches(0.42), sz=15, bold=True, clr=ACC)
    right_lines = [
        "Every action logged with:",
        "• UTC timestamp",
        "• Action type: criteria_extracted / bidder_processed",
        "  criterion_evaluated / human_review_action",
        "  vision_ocr_invoked / precomputed_fallback_used",
        "• Model version & Actor (system / officer)",
        "• Full payload JSON — Exportable as CSV",
    ]
    for i, line in enumerate(right_lines):
        c2 = TXT2 if i == 0 else TXT
        txt(slide, line, rx + Inches(0.2), Inches(1.68) + Inches(0.42 * i),
            Inches(5.5), Inches(0.38), sz=13, clr=c2)

    rect(slide, Inches(0.6), Inches(5.38), Inches(12.13), Inches(1.22),
         RGBColor(0xFE, 0xF3, 0xC7), AMB, Pt(2))
    txt(slide, "The Safety Rule:",
        Inches(0.8), Inches(5.45), Inches(5), Inches(0.38),
        sz=15, bold=True, clr=AMB)
    txt(slide, ("If combined confidence is 0.55–0.80 AND verdict is not_eligible, "
                "the verdict is automatically downgraded to needs_review. "
                "A bidder is NEVER silently disqualified at medium confidence."),
        Inches(0.8), Inches(5.88), Inches(12), Inches(0.65),
        sz=13, clr=TXT)


# ── Slide 7 ─────────────────────────────────────────────────────────────────
def s7(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Demo: Three Bidders, Three Outcomes")
    mpl_style()

    # Verdicts matrix: 3 bidders × 5 criteria
    criteria = ["C1\nTurnover", "C2\nProjects", "C3\nGST", "C4\nISO", "C5\nExperience"]
    bidders = ["Bidder A\n(Apex)", "Bidder B\n(BuildRight)", "Bidder C\n(Shree)"]
    verdicts = [
        ["E", "E", "E", "E", "E"],
        ["N", "E", "E", "E", "E"],
        ["R", "R", "E", "E", "E"],
    ]
    color_map = {"E": HEX["grn"], "N": HEX["red"], "R": HEX["amb"]}
    label_map = {"E": "E", "N": "N", "R": "R"}

    fig, ax = plt.subplots(figsize=(5.2, 3.2))
    for r, row in enumerate(verdicts):
        for col, v in enumerate(row):
            c_ = color_map[v]
            ax.add_patch(mpatches.FancyBboxPatch(
                (col + 0.05, r + 0.05), 0.9, 0.9,
                boxstyle="round,pad=0.05", facecolor=c_, edgecolor="white", lw=2))
            ax.text(col + 0.5, r + 0.5, label_map[v],
                    ha="center", va="center", fontsize=16,
                    fontweight="bold", color="white")
    ax.set_xlim(0, 5)
    ax.set_ylim(0, 3)
    ax.set_xticks([i + 0.5 for i in range(5)])
    ax.set_xticklabels(criteria, fontsize=9)
    ax.set_yticks([i + 0.5 for i in range(3)])
    ax.set_yticklabels(bidders, fontsize=9)
    ax.set_title("Verdict Matrix (E=Eligible, N=Not Eligible, R=Review)",
                 color=HEX["pri"], fontsize=10, pad=6)
    ax.grid(False)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["bottom"].set_visible(False)
    ax.spines["left"].set_visible(False)
    png_to_slide(slide, fig, Inches(0.6), Inches(1.0), Inches(5.8), Inches(4.0))

    bidders_info = [
        ("Bidder A — Apex Constructions", GRN, "ELIGIBLE",
         ["C1 INR 6.37 Cr avg — PASS",
          "C2 5 completed projects — PASS",
          "C3 GST active — PASS",
          "C4 ISO 9001:2015 valid — PASS",
          "Confidence ≥ 93% all criteria"]),
        ("Bidder B — BuildRight Enterprises", RED_C, "NOT ELIGIBLE",
         ["C1 INR 1.5 Cr avg — FAIL",
          "Below required INR 5 Cr",
          "C2–C4: All pass",
          "Disqualified, high conf (95%)"]),
        ("Bidder C — Shree Constructions", AMB, "NEEDS REVIEW",
         ["C1 Blurry scan: Vision LLM",
          "INR 5.4 Cr, conf 0.58 → review",
          "C2: 3 projects (borderline)",
          "C3–C4: Pass"]),
    ]
    cw = Inches(2.2)
    for i, (name, color, verdict, bullets) in enumerate(bidders_info):
        sx = Inches(6.9) + i * Inches(2.15)
        rect(slide, sx, Inches(1.0), cw, Inches(4.0),
             RGBColor(0xF8, 0xF9, 0xFF), color, Pt(2))
        rect(slide, sx, Inches(1.0), cw, Inches(0.08), color)
        txt(slide, name, sx + Inches(0.1), Inches(1.1), cw - Inches(0.2), Inches(0.45),
            sz=11, bold=True, clr=PRI)
        rect(slide, sx + Inches(0.1), Inches(1.62), cw - Inches(0.2), Inches(0.32), color)
        txt(slide, verdict, sx + Inches(0.1), Inches(1.62), cw - Inches(0.2), Inches(0.32),
            sz=11, bold=True, clr=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)
        for j, b in enumerate(bullets):
            txt(slide, b, sx + Inches(0.1), Inches(2.05) + Inches(0.38 * j),
                cw - Inches(0.2), Inches(0.34), sz=11, clr=TXT)

    # Metric strip
    metrics = [
        ("Criteria extracted", "5"),
        ("Bidder docs processed", "15"),
        ("LLM evaluation calls", "15"),
        ("Vision OCR invocations", "1"),
        ("Human review items", "1"),
        ("Total audit entries", "20+"),
    ]
    rect(slide, Inches(0.6), Inches(5.2), Inches(12.13), Inches(0.9),
         RGBColor(0xEE, 0xF2, 0xFF), ACC, Pt(1))
    mw = Inches(2.0)
    for i, (lbl, val) in enumerate(metrics):
        mx = Inches(0.7) + i * Inches(2.0)
        txt(slide, val, mx, Inches(5.27), mw, Inches(0.38),
            sz=20, bold=True, clr=ACC, align=PP_ALIGN.CENTER)
        txt(slide, lbl, mx, Inches(5.7), mw, Inches(0.32),
            sz=11, clr=TXT2, align=PP_ALIGN.CENTER)


# ── Slide 8 ─────────────────────────────────────────────────────────────────
def s8(prs):
    slide = blank(prs)
    bg(slide)
    heading(slide, "Stack, Impact & What's Next")
    stack = [
        ("UI & orchestration", "Streamlit 1.39"),
        ("LLM", "DeepSeek API (OpenAI-compatible)"),
        ("OCR Tier 1", "PyMuPDF 1.24"),
        ("OCR Tier 2", "Tesseract"),
        ("OCR Tier 3", "DeepSeek Vision LLM"),
        ("Semantic retrieval", "sentence-transformers all-MiniLM-L6-v2"),
        ("Data validation", "Pydantic v2"),
        ("Audit log", "SQLite"),
        ("Deployment", "Streamlit Cloud / HuggingFace Spaces"),
    ]
    rect(slide, Inches(0.6), Inches(1.05), Inches(6.1), Inches(0.4),
         RGBColor(0xEE, 0xF2, 0xFF), ACC, Pt(1))
    txt(slide, "Component", Inches(0.7), Inches(1.1), Inches(2.7), Inches(0.3),
        sz=13, bold=True, clr=ACC)
    txt(slide, "Technology", Inches(3.55), Inches(1.1), Inches(2.9), Inches(0.3),
        sz=13, bold=True, clr=ACC)
    for i, (comp, tech) in enumerate(stack):
        ry = Inches(1.45) + Inches(0.38 * i)
        bg_r = RGBColor(0xF8, 0xF9, 0xFF) if i % 2 == 0 else BG
        rect(slide, Inches(0.6), ry, Inches(6.1), Inches(0.38), bg_r)
        txt(slide, comp, Inches(0.7), ry + Inches(0.06), Inches(2.7), Inches(0.3),
            sz=12, clr=TXT2)
        txt(slide, tech, Inches(3.55), ry + Inches(0.06), Inches(2.9), Inches(0.3),
            sz=12, clr=TXT)

    txt(slide, "What's Next", Inches(7.3), Inches(1.05), Inches(5.6), Inches(0.42),
        sz=16, bold=True, clr=ACC)
    future = [
        "Multi-tender workspace — same bidder pool, multiple tenders",
        "GeM portal API integration — live tender ingestion",
        "Automated bidder ranking with weighted scoring",
        "LayoutLM for complex financial tables in scanned statements",
        "Multi-evaluator workflow with role-based approval",
        "Review queue email/SMS notifications",
        "Audit PDF export for procurement oversight submissions",
    ]
    for i, f in enumerate(future):
        txt(slide, f"• {f}", Inches(7.3), Inches(1.55) + Inches(0.47 * i),
            Inches(5.8), Inches(0.42), sz=13, clr=TXT)

    rect(slide, Inches(0.6), Inches(6.2), Inches(12.13), Inches(0.92),
         RGBColor(0xEE, 0xF2, 0xFF), ACC, Pt(2))
    txt(slide, ("3–5 days → minutes.  Every verdict traceable to a document, page, and model version."
                "  Built in one hackathon session. Deployable today."),
        Inches(0.8), Inches(6.3), Inches(11.9), Inches(0.75),
        sz=14, bold=True, clr=ACC, align=PP_ALIGN.CENTER)


def main():
    os.makedirs("deck", exist_ok=True)
    prs = new_prs()
    s1(prs); s2(prs); s3(prs); s4(prs)
    s5(prs); s6(prs); s7(prs); s8(prs)
    out = "deck/TenderIQ_v5_data_forward.pptx"
    prs.save(out)
    print(f"OK: Saved {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
